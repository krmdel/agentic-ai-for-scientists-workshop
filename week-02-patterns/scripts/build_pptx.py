"""Build a NATIVE, fully-editable Week 2 PowerPoint (python-pptx).

Unlike the Marp export (one image per slide) or Marp --pptx-editable (LibreOffice
fragments every span into its own box), this emits real PowerPoint:
  - one text frame per region; each bullet is ONE paragraph (bold/code as inline runs)
  - native editable tables
  - brand styling (TUMCREATE blue, logo, title building); cards/code = simple native shapes

Run:  .venv/bin/python3.11 scripts/build_pptx.py   ->  ../slides.pptx
"""
from __future__ import annotations
import math, re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
LOGO = str(ROOT / "assets/logos/tumcreate.png")
BUILDING = str(ROOT / "assets/logos/tumcreate-intro.png")
TXFIG = str(ROOT / "assets/figures/transformer-architecture-vaswani-2017.png")
OUT = ROOT / "slides.pptx"

# ---- palette ----
INK   = RGBColor.from_string("1A1A1A")
HEAD  = RGBColor.from_string("0A0A0A")
BLUE  = RGBColor.from_string("1F4F8B")
GREEN = RGBColor.from_string("4A7D3F")
CORAL = RGBColor.from_string("C8553D")
MUTED = RGBColor.from_string("6A6A6A")
CARDBG= RGBColor.from_string("F7F8FA")
CODEBG= RGBColor.from_string("F3F4F6")
RUNBG = RGBColor.from_string("EAF4EA")
RUNFG = RGBColor.from_string("2F6D2A")
WHITE = RGBColor.from_string("FFFFFF")
ACCENT = {"blue": BLUE, "green": GREEN, "coral": CORAL}

SANS = "Helvetica Neue"
MONO = "Menlo"

# Week 2 date — footer bottom-left, matching Week 1's DD/MM/YYYY notation.
TALK_DATE = "03/06/2026"

SW, SH = Inches(13.333), Inches(7.5)
ML = 0.62
FULL_W = 12.09
COLS_EQ = [(0.62, 5.75), (6.96, 5.75)]
COLS_WIDE = [(0.62, 6.78), (7.74, 4.97)]
BODY_TOP = 1.66
BODY_BOT = 6.62

INLINE = re.compile(r'(\*\*.+?\*\*|`[^`]+?`|\*[^*]+?\*)')


def runs_of(text):
    out = []
    for seg in INLINE.split(text):
        if not seg:
            continue
        if seg.startswith("**") and seg.endswith("**"):
            out.append((seg[2:-2], True, False, False))
        elif seg.startswith("`") and seg.endswith("`"):
            out.append((seg[1:-1], False, False, True))
        elif seg.startswith("*") and seg.endswith("*") and len(seg) > 2:
            out.append((seg[1:-1], False, True, False))
        else:
            out.append((seg, False, False, False))
    return out


def plain(text):
    return text.replace("**", "").replace("`", "").replace("*", "")


def add_runs(p, text, size, color, bold=False):
    for t, b, i, m in runs_of(text):
        r = p.add_run(); r.text = t
        f = r.font
        f.size = Pt(size); f.bold = bold or b; f.italic = i
        f.name = MONO if m else SANS
        f.color.rgb = color


def _tf(shape):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    return tf


def box(slide, x, y, w, h):
    return slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))


def set_indent(p, marL=0.30):
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(int(marL * 914400)))
    pPr.set("indent", str(int(-marL * 914400)))


# ---------- height estimation (inches) ----------
def cpl(w_in, pt):
    return max(8, int(w_in * 145 / pt))


def lh(pt):
    return pt / 72.0 * 1.30


def lines(text, w_in, pt):
    t = plain(text)
    n = 0
    for seg in t.split("\n"):
        n += max(1, math.ceil(len(seg) / cpl(w_in, pt)))
    return n


def est_h(block, w_in):
    k = block[0]
    if k == "h":   return lh(16) * lines(block[1], w_in, 16) + 0.10
    if k == "h3":  return lh(19) * lines(block[1], w_in, 19) + 0.10
    if k == "p":   return lh(16) * lines(block[1], w_in, 16) + 0.12
    if k in ("b", "n"):
        bw = w_in - 0.30
        h = 0.04
        for it in block[1]:
            h += lh(16) * lines(it, bw, 16) + 0.07
        return h + 0.06
    if k == "code":
        nl = block[1].count("\n") + 1
        return nl * lh(12) + 0.26
    if k == "card":
        bw = w_in - 0.5
        bh = sum(lh(14) * lines(par, bw, 14) + 0.06 for par in block[3].split("||"))
        return 0.30 + bh + 0.34
    if k == "muted": return lh(12) * lines(block[1], w_in, 12) + 0.10
    if k == "table": return len(block[2]) * 0.36 + 0.18 + (0.36)
    if k == "run":   return 0.52
    if k == "dgm":   return DGM_H[block[1]] + 0.18
    return 0.4


# ---------- block renderers (return bottom y) ----------
def r_text(slide, x, y, w, block, size, color, bold=False, space=0.12, align=PP_ALIGN.LEFT):
    h = est_h(block, w)
    tb = box(slide, x, y, w, h); tf = _tf(tb)
    p = tf.paragraphs[0]; p.alignment = align
    add_runs(p, block[1], size, color, bold)
    return y + h - 0.02 + space


def r_list(slide, x, y, w, items, numbered=False):
    h = est_h(("b" if not numbered else "n", items), w)
    tb = box(slide, x, y, w, h); tf = _tf(tb)
    for idx, it in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(5); set_indent(p, 0.30)
        g = p.add_run(); g.text = (f"{idx+1}.  " if numbered else "▪  ")
        g.font.size = Pt(16); g.font.name = SANS
        g.font.color.rgb = MUTED if not numbered else BLUE
        if numbered: g.font.bold = True
        add_runs(p, it, 16, INK)
    return y + h + 0.08


def r_code(slide, x, y, w, text):
    nl = text.count("\n") + 1
    h = nl * lh(12) + 0.24
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = CODEBG
    sh.line.fill.background(); sh.shadow.inherit = False
    tf = _tf(sh); tf.margin_left = Inches(0.14); tf.margin_top = Inches(0.10)
    tf.margin_right = Inches(0.10); tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.0; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = ln if ln else " "
        r.font.size = Pt(12); r.font.name = MONO; r.font.color.rgb = INK
    return y + h + 0.12


def r_card(slide, x, y, w, label, accent, body):
    bw = w - 0.5
    bh = sum(lh(14) * lines(par, bw, 14) + 0.06 for par in body.split("||"))
    h = 0.30 + bh + 0.30
    col = ACCENT[accent]
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.adjustments[0] = 0.04
    rect.fill.solid(); rect.fill.fore_color.rgb = CARDBG
    rect.line.fill.background(); rect.shadow.inherit = False
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background(); bar.shadow.inherit = False
    tf = _tf(rect); tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.13); tf.margin_bottom = Inches(0.10)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    lr = p.add_run(); lr.text = label.upper()
    lr.font.size = Pt(10.5); lr.font.bold = True; lr.font.name = SANS; lr.font.color.rgb = col
    for par in body.split("||"):
        bp = tf.add_paragraph(); bp.space_before = Pt(4); bp.alignment = PP_ALIGN.LEFT
        add_runs(bp, par, 14, INK)
    return y + h + 0.14


def r_table(slide, x, y, w, headers, rows):
    nr, nc = len(rows) + 1, len(headers)
    h = nr * 0.36
    gt = slide.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(h)).table
    gt.first_row = True; gt.horz_banding = False
    for j, htext in enumerate(headers):
        c = gt.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = WHITE
        c.margin_left = Inches(0.08); c.margin_top = Inches(0.03); c.margin_bottom = Inches(0.03)
        p = c.text_frame.paragraphs[0]; add_runs(p, htext, 12, HEAD, bold=True)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = gt.cell(i, j); c.fill.solid(); c.fill.fore_color.rgb = WHITE
            c.margin_left = Inches(0.08); c.margin_top = Inches(0.03); c.margin_bottom = Inches(0.03)
            p = c.text_frame.paragraphs[0]; add_runs(p, val, 12, INK)
    return y + h + 0.14


def r_run(slide, x, y, w, text):
    pw = min(w, 7.6)
    nl = max(1, math.ceil(len(text) / cpl(pw - 0.3, 13)))
    h = max(0.40, nl * lh(13) + 0.20)
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(pw), Inches(h))
    sh.adjustments[0] = 0.3
    sh.fill.solid(); sh.fill.fore_color.rgb = RUNBG; sh.line.fill.background(); sh.shadow.inherit = False
    tf = _tf(sh); tf.margin_left = Inches(0.16); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT; r = p.add_run(); r.text = text
    r.font.size = Pt(13); r.font.bold = True; r.font.name = SANS; r.font.color.rgb = RUNFG
    return y + h + 0.12


def render_blocks(slide, x, y, w, blocks):
    for blk in blocks:
        k = blk[0]
        if k == "h":   y = r_text(slide, x, y, w, blk, 16, HEAD, bold=True, space=0.10)
        elif k == "h3": y = r_text(slide, x, y, w, blk, 19, HEAD, bold=True, space=0.08)
        elif k == "p":  y = r_text(slide, x, y, w, blk, 16, INK)
        elif k == "b":  y = r_list(slide, x, y, w, blk[1])
        elif k == "n":  y = r_list(slide, x, y, w, blk[1], numbered=True)
        elif k == "code": y = r_code(slide, x, y, w, blk[1])
        elif k == "card": y = r_card(slide, x, y, w, blk[1], blk[2], blk[3])
        elif k == "muted": y = r_text(slide, x, y, w, blk, 12, MUTED, space=0.10)
        elif k == "table": y = r_table(slide, x, y, w, blk[1], blk[2])
        elif k == "run": y = r_run(slide, x, y, w, blk[1])
        elif k == "dgm": DGM[blk[1]](slide, x, y, w); y += DGM_H[blk[1]] + 0.18
    return y


# ---------- chrome ----------
def add_logo(slide):
    slide.shapes.add_picture(LOGO, Inches(11.05), Inches(0.36), width=Inches(2.0))


def add_footer(slide, page):
    lb = box(slide, ML, 6.96, 4.0, 0.3); p = _tf(lb).paragraphs[0]
    r = p.add_run(); r.text = TALK_DATE
    r.font.size = Pt(10); r.font.name = SANS; r.font.color.rgb = MUTED
    cb = box(slide, 3.5, 6.96, 6.33, 0.3); p = _tf(cb).paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "kerem.delikoyun@tum-create.edu.sg   |   Agentic AI Workshop W2"
    r.font.size = Pt(10); r.font.name = SANS; r.font.color.rgb = MUTED
    nb = box(slide, 12.2, 6.96, 0.7, 0.3); p = _tf(nb).paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = str(page)
    r.font.size = Pt(10); r.font.name = SANS; r.font.color.rgb = MUTED


def add_title(slide, text):
    tb = box(slide, ML, 0.78, 11.6, 0.9); tf = _tf(tb)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(29); r.font.bold = True; r.font.name = SANS; r.font.color.rgb = HEAD


# ---------- slide builders ----------
def build_title(prs, title, lede, byline, building=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(s)
    tb = box(s, ML, 2.35, 9.5, 1.4); p = _tf(tb).paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(46); r.font.bold = True; r.font.name = SANS; r.font.color.rgb = HEAD
    lb = box(s, ML, 3.95, 9.6, 1.0); p = _tf(lb).paragraphs[0]
    add_runs(p, lede, 20, RGBColor.from_string("444444"))
    bb = box(s, ML, 5.15, 9.6, 0.5); p = _tf(bb).paragraphs[0]
    r = p.add_run(); r.text = byline
    r.font.size = Pt(14); r.font.name = SANS; r.font.color.rgb = RGBColor.from_string("666666")
    if building:
        s.shapes.add_picture(BUILDING, Inches(9.55), Inches(4.55), width=Inches(3.2))
    return s


def build_opener(prs, page, chapter, title, blurb):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(s); add_footer(s, page)
    cb = box(s, ML, 2.55, 11.0, 0.4); p = _tf(cb).paragraphs[0]
    r = p.add_run(); r.text = chapter.upper()
    r.font.size = Pt(13); r.font.bold = True; r.font.name = SANS; r.font.color.rgb = BLUE
    tb = box(s, ML, 3.0, 11.6, 1.0); p = _tf(tb).paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(40); r.font.bold = True; r.font.name = SANS; r.font.color.rgb = HEAD
    bb = box(s, ML, 4.15, 10.0, 1.2); p = _tf(bb).paragraphs[0]
    add_runs(bb.text_frame.paragraphs[0], blurb, 18, RGBColor.from_string("555555"))
    return s


def build_content(prs, page, slide):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(s); add_footer(s, page); add_title(s, slide["title"])
    cols = COLS_WIDE if slide.get("wide") else COLS_EQ
    y = BODY_TOP
    for row in slide["body"]:
        if row[0] == "full":
            render_blocks(s, ML, y, FULL_W, row[1])
            # advance y by estimated full height
            yy = y
            for b in row[1]:
                yy += est_h(b, FULL_W)
            y = yy + 0.06
        else:  # cols
            (lx, lw), (rx, rw) = cols
            render_blocks(s, lx, y, lw, row[1])
            render_blocks(s, rx, y, rw, row[2])
            lh_ = sum(est_h(b, lw) for b in row[1])
            rh_ = sum(est_h(b, rw) for b in row[2])
            y += max(lh_, rh_) + 0.10
    if slide.get("image"):
        path, ix, iy, iw = slide["image"]
        s.shapes.add_picture(path, Inches(ix), Inches(iy), width=Inches(iw))
    if slide.get("foot"):
        fb = slide["foot"]
        render_blocks(s, ML, max(y, 6.15), FULL_W, [fb])
    return s


# ============================================================================
# NATIVE DIAGRAMS  (replace the box-drawing ASCII)
# ============================================================================
def _noshadow(sh):
    sh.shadow.inherit = False


def nodebox(slide, x, y, w, h, title, accent="blue", sub=None, solid=False, size=14, rounded=True):
    col = ACCENT[accent]
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        sh.adjustments[0] = 0.12
    sh.fill.solid(); sh.fill.fore_color.rgb = col if solid else CARDBG
    sh.line.color.rgb = col; sh.line.width = Pt(1.25); _noshadow(sh)
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    r.font.size = Pt(size); r.font.bold = True; r.font.name = SANS
    r.font.color.rgb = WHITE if solid else col
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(2)
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(10.5); r2.font.name = SANS
        r2.font.color.rgb = (WHITE if solid else MUTED)
    return sh


def plainbox(slide, x, y, w, h, fill=CARDBG, line=None, rounded=True):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        sh.adjustments[0] = 0.10
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is not None:
        sh.line.color.rgb = line; sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    _noshadow(sh); return sh


def line_seg(slide, x1, y1, x2, y2, color=MUTED, width=1.5):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(width); _noshadow(cn)
    return cn


def arrow(slide, x1, y1, x2, y2, color=MUTED, width=1.5, dbl=False):
    cn = line_seg(slide, x1, y1, x2, y2, color, width)
    ln = cn.line._get_or_add_ln()
    if dbl:
        he = ln.makeelement(qn('a:headEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}); ln.append(he)
    te = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}); ln.append(te)
    return cn


def caption(slide, x, y, w, text, size=11, color=MUTED, align=PP_ALIGN.CENTER, bold=False, italic=False, mono=False):
    tb = box(slide, x, y, w, 0.32); tf = _tf(tb); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.name = MONO if mono else SANS; r.font.color.rgb = color
    return tb


def _bar_row(slide, x, y, w, frac, name, val, color, lbl_w=1.55):
    caption(slide, x, y - 0.01, lbl_w, name, size=12, color=INK, align=PP_ALIGN.RIGHT)
    tx = x + lbl_w + 0.12
    tw = w - lbl_w - 0.12 - 0.62
    plainbox(slide, tx, y + 0.05, tw, 0.22, fill=CODEBG, rounded=False)
    plainbox(slide, tx, y + 0.05, max(0.10, tw * frac), 0.22, fill=color, rounded=False)
    caption(slide, tx + tw + 0.06, y - 0.01, 0.6, val, size=12, color=MUTED, align=PP_ALIGN.LEFT)
    return y + 0.40


# ---- the diagrams (each draws into x,y with given width w) ----
def dgm_anatomy(slide, x, y, w):
    cx = x + w / 2
    bw, bh = 3.2, 0.74
    nodebox(slide, cx - bw / 2, y, bw, bh, "LLM", accent="blue", sub="the brain", solid=True, size=17)
    cw, ch = 3.5, 1.30
    gap = (w - 3 * cw) / 2
    xs = [x, x + cw + gap, x + 2 * (cw + gap)]
    cy = y + bh + 1.0
    specs = [("Tools", "green", "APIs · DBs · code · functions"),
             ("Memory", "blue", "context window + RAG"),
             ("Control loop", "coral", "you write this!")]
    for tx, (ti, acc, sub) in zip(xs, specs):
        nodebox(slide, tx, cy, cw, ch, ti, accent=acc, sub=sub, size=16)
    busy = y + bh + 0.50
    line_seg(slide, cx, y + bh, cx, busy)
    centers = [tx + cw / 2 for tx in xs]
    line_seg(slide, centers[0], busy, centers[-1], busy)
    for c in centers:
        arrow(slide, c, busy, c, cy)
    return y + bh + 1.0 + ch


def dgm_ladder(slide, x, y, w):
    rows = [("L3", "Multi-agent systems", "specialists + orchestration", "coral", "Week 3"),
            ("L2", "+ Reasoning", "CoT · ToT · ReAct", "green", "today"),
            ("L1", "+ Tool use", "the model acts on the world", "green", "today"),
            ("L0", "Bare LLM", "one prompt, one answer", "blue", "today")]
    rh, gap = 0.70, 0.12
    chip = 0.92
    tag_w = 1.15
    body_x = x + chip + 0.14
    body_w = w - chip - 0.14 - tag_w - 0.12
    for i, (rung, title, desc, acc, tag) in enumerate(rows):
        ry = y + i * (rh + gap)
        nodebox(slide, x, ry, chip, rh, rung, accent=acc, solid=True, size=17)
        plainbox(slide, body_x, ry, body_w, rh, fill=CARDBG)
        tb = box(slide, body_x + 0.22, ry, body_w - 0.4, rh); tf = _tf(tb)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = title + "    "
        r.font.size = Pt(15); r.font.bold = True; r.font.name = SANS; r.font.color.rgb = HEAD
        r2 = p.add_run(); r2.text = desc
        r2.font.size = Pt(13); r2.font.name = SANS; r2.font.color.rgb = MUTED
        tcol = CORAL if tag == "Week 3" else GREEN
        pill = plainbox(slide, x + w - tag_w, ry + rh / 2 - 0.18, tag_w, 0.36, fill=CARDBG, line=tcol)
        tf2 = pill.text_frame; tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = tf2.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        rr = pp.add_run(); rr.text = ("→ " + tag) if tag == "Week 3" else tag
        rr.font.size = Pt(11); rr.font.bold = True; rr.font.name = SANS; rr.font.color.rgb = tcol
    bottom = y + 4 * rh + 3 * gap
    return bottom


def dgm_l1loop(slide, x, y, w):
    steps = ["LLM reads the question + tool list",
             "LLM requests a call:  calculator(17×25)",
             "YOUR code runs the function",
             "feed the result back to the LLM",
             "repeat until the LLM is done"]
    sh_h, gap = 0.46, 0.18
    bw = w - 1.0
    for i, s in enumerate(steps):
        sy = y + i * (sh_h + gap)
        acc = "green" if i == 2 else "blue"
        nodebox(slide, x, sy, 0.46, sh_h, str(i + 1), accent=acc, solid=True, size=14)
        b = plainbox(slide, x + 0.6, sy, bw - 0.6, sh_h, fill=CARDBG)
        tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.14)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        add_runs(p, s, 12.5, INK)
        if i < len(steps) - 1:
            arrow(slide, x + 0.23, sy + sh_h, x + 0.23, sy + sh_h + gap)
    top = y + 0.23
    bot = y + 4 * (sh_h + gap) + 0.23
    rx = x + w - 0.30
    line_seg(slide, x + bw + 0.02, bot - 0.1, rx, bot - 0.1, color=GREEN, width=1.5)
    line_seg(slide, rx, bot - 0.1, rx, top, color=GREEN, width=1.5)
    arrow(slide, rx, top, x + bw + 0.02, top, color=GREEN, width=1.5)
    caption(slide, rx - 0.7, (top + bot) / 2 - 0.16, 0.95, "loop", size=11, color=GREEN, bold=True, align=PP_ALIGN.RIGHT)
    return bot + 0.1


def dgm_tot(slide, x, y, w):
    cx = x + w / 2
    nodebox(slide, cx - 1.1, y, 2.2, 0.5, "problem", accent="blue", size=14)
    cw = 1.75
    gap = (w - 3 * cw) / 2
    xs = [x, x + cw + gap, x + 2 * (cw + gap)]
    cy = y + 1.05
    evals = [("✓ sure", "green"), ("~ maybe", "blue"), ("✗ dead", "coral")]
    for tx, (tag, acc) in zip(xs, evals):
        c = tx + cw / 2
        arrow(slide, cx, y + 0.5, c, cy)
        nodebox(slide, tx, cy, cw, 0.5, "branch", accent=acc, sub=tag, size=12.5)
    fy = cy + 0.95
    # expand the best, prune the dead
    arrow(slide, xs[0] + cw / 2, cy + 0.5, xs[0] + cw / 2, fy, color=GREEN)
    caption(slide, xs[0] - 0.1, fy, cw + 0.2, "expand", size=11, color=GREEN, bold=True)
    caption(slide, xs[2] - 0.1, fy, cw + 0.2, "prune", size=11, color=CORAL, bold=True)
    return fy + 0.34


def dgm_attn(slide, x, y, w):
    b = plainbox(slide, x, y, w, 0.56, fill=CARDBG)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    tf.margin_left = Inches(0.16); tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    for seg, col, bold in [("The scientist read the paper because ", INK, False),
                           ("it", BLUE, True), (" was ___", INK, False)]:
        r = p.add_run(); r.text = seg
        r.font.size = Pt(14); r.font.bold = bold; r.font.name = SANS; r.font.color.rgb = col
    yy = y + 0.74
    caption(slide, x, yy, w, "what  “it”  attends to — by learned relevance:", size=11.5, color=MUTED, align=PP_ALIGN.LEFT)
    yy += 0.34
    yy = _bar_row(slide, x, yy, w, 0.70, "paper", "0.70", BLUE)
    yy = _bar_row(slide, x, yy, w, 0.20, "scientist", "0.20", RGBColor.from_string("9DB7D6"))
    return yy


def dgm_nexttok(slide, x, y, w):
    b = plainbox(slide, x, y, w * 0.92, 0.5, fill=CARDBG)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.16)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "The capital of France is "
    r.font.size = Pt(13.5); r.font.name = SANS; r.font.color.rgb = INK
    r2 = p.add_run(); r2.text = "____"
    r2.font.size = Pt(13.5); r2.font.bold = True; r2.font.name = MONO; r2.font.color.rgb = BLUE
    cx = x + 0.6
    arrow(slide, cx, y + 0.5, cx, y + 0.78)
    caption(slide, cx + 0.16, y + 0.5, 3.0, "P(next token)", size=10.5, color=MUTED, align=PP_ALIGN.LEFT)
    yy = y + 0.86
    yy = _bar_row(slide, x, yy, w, 0.91, "“ Paris”", "0.91", BLUE)
    yy = _bar_row(slide, x, yy, w, 0.05, "“ a”", "0.02", RGBColor.from_string("9DB7D6"))
    yy = _bar_row(slide, x, yy, w, 0.03, "“ home”", "0.01", RGBColor.from_string("C7D3E3"))
    return yy


def _flow(slide, x, y, w, items, accents):
    n = len(items)
    gap = 0.30
    bw = (w - gap * (n - 1)) / n
    h = 0.62
    for i, (it, acc) in enumerate(zip(items, accents)):
        bx = x + i * (bw + gap)
        nodebox(slide, bx, y, bw, h, it, accent=acc, size=12)
        if i < n - 1:
            arrow(slide, bx + bw, y + h / 2, bx + bw + gap, y + h / 2)
    return y + h


def dgm_ragdiag(slide, x, y, w):
    caption(slide, x, y, 3.0, "INDEX  (once)", size=10.5, color=MUTED, align=PP_ALIGN.LEFT, bold=True)
    y1 = y + 0.30
    by = _flow(slide, x, y1, w, ["PDFs", "text", "chunks", "embeddings", "vector store"],
               ["blue"] * 4 + ["green"])
    caption(slide, x, by + 0.18, 3.0, "QUERY  (every question)", size=10.5, color=MUTED, align=PP_ALIGN.LEFT, bold=True)
    y2 = by + 0.48
    by2 = _flow(slide, x, y2, w, ["question", "embed", "top-k chunks", "prompt + chunks", "LLM", "answer + citations"],
                ["blue", "blue", "green", "blue", "blue", "coral"])
    return by2


def dgm_w3node(slide, x, y, w):
    qy = y + 0.6
    nodebox(slide, x, qy, 1.7, 0.6, "user query", accent="blue", size=12)
    ox = x + 2.3
    ow = w - 2.3
    cont = plainbox(slide, ox, y, ow, 2.0, fill=RGBColor.from_string("F4F7FB"), line=BLUE)
    caption(slide, ox + 0.2, y + 0.12, ow - 0.4, "LangGraph orchestrator", size=12.5, color=BLUE, align=PP_ALIGN.LEFT, bold=True)
    arrow(slide, x + 1.7, qy + 0.3, ox, qy + 0.3, color=BLUE)
    iw, ih = 2.4, 0.62
    igap = (ow - 0.8 - 3 * iw) / 2
    ix0 = ox + 0.4
    iy = y + 0.72
    names = [("Planner", "green"), ("Retriever", "blue"), ("Writer", "coral")]
    ixs = [ix0 + k * (iw + igap) for k in range(3)]
    for (ix, (nm, acc)) in zip(ixs, names):
        nodebox(slide, ix, iy, iw, ih, nm, accent=acc, size=13)
    arrow(slide, ixs[0] + iw, iy + ih / 2, ixs[1], iy + ih / 2)
    arrow(slide, ixs[1] + iw, iy + ih / 2, ixs[2], iy + ih / 2)
    caption(slide, ixs[1] - 0.3, iy + ih + 0.16, iw + 0.6,
            "hybrid_retriever()  —  straight from Notebook 03", size=11, color=MUTED)
    return y + 2.0


DGM = {"anatomy": dgm_anatomy, "ladder": dgm_ladder, "l1loop": dgm_l1loop, "tot": dgm_tot,
       "attn": dgm_attn, "nexttok": dgm_nexttok, "ragdiag": dgm_ragdiag, "w3node": dgm_w3node}
DGM_H = {"anatomy": 3.05, "ladder": 3.22, "l1loop": 3.34, "tot": 2.40,
         "attn": 2.10, "nexttok": 2.32, "ragdiag": 2.40, "w3node": 2.10}


# ============================================================================
# DECK CONTENT
# ============================================================================
LADDER = (
    "  L3   Multi-agent systems        specialists + orchestration   ->  Week 3\n"
    "       ---------------------------------------------------------------\n"
    "  L2   + Reasoning                CoT . ToT . ReAct             <- today\n"
    "  L1   + Tool use                 the model acts on the world   <- today\n"
    "  L0   Bare LLM                   one prompt, one answer        <- today"
)
ANATOMY = (
    "                  +--------------------+\n"
    "                  |   LLM (the brain)  |\n"
    "                  +---------+----------+\n"
    "                            |\n"
    "   +------------------------+------------------------+\n"
    "   |                        |                        |\n"
    "+--v---+               +----v-----+            +-----v-----+\n"
    "| Tools|               |  Memory  |            |  Control  |\n"
    "| APIs |               | (context |            |   Loop    |\n"
    "| DBs  |               |  + RAG)  |            | (you write|\n"
    "| Code |               +----------+            |  this!)   |\n"
    "+------+                                       +-----------+"
)
NEXTTOK = (
    "The capital of France is  ____\n"
    "                          |\n"
    "                          v\n"
    "      P(next token)  ->  \" Paris\"   0.91\n"
    "                        \" a\"       0.02\n"
    "                        \" home\"    0.01\n"
    "                        ..."
)
ATTN = (
    "  The  scientist  read  the  paper  because  it  was  ___\n"
    "\n"
    "                                            |\n"
    "                +---------------------------+\n"
    "          \"it\" attends to -> \"paper\"  (0.7)\n"
    "                          -> \"scientist\" (0.2)"
)
L1LOOP = (
    "+--------------------------------------+\n"
    "|  1. LLM reads question + tool list   |\n"
    "|  2. LLM: \"call calculator(17*25)\"    |\n"
    "|  3. YOU run the function             |\n"
    "|  4. feed the result back             |\n"
    "|  5. loop until the LLM is done       |\n"
    "+--------------------------------------+"
)
COTBLK = (
    "Q: A juggler has 16 balls. Half are golf balls,\n"
    "   and half of those are blue. How many blue\n"
    "   golf balls?\n"
    "\n"
    "DIRECT:           4        (right or wrong, no trace)\n"
    "\n"
    "STEP BY STEP:\n"
    "  16 balls -> half golf      = 8 golf balls\n"
    "  8 golf  -> half blue       = 4 blue golf balls\n"
    "  Answer: 4                  (visible, checkable)"
)
TOTBLK = (
    "                  problem\n"
    "                  /  |  \\\n"
    "            branch  branch  branch     <- PROPOSE several thoughts\n"
    "              |       |       |\n"
    "            eval    eval    eval        <- EVALUATE: sure / maybe / dead\n"
    "              |               x\n"
    "            expand          prune\n"
    "              |\n"
    "            ...continue the promising branches"
)
REACTPROMPT = (
    "Answer the question. Tools: {tools}\n"
    "\n"
    "Question: the input question\n"
    "Thought:  reasoning about what to do next\n"
    "Action:   one of [{tool_names}]\n"
    "Action Input: the input to the action\n"
    "Observation: the result of the action\n"
    "... (repeat) ...\n"
    "Thought:  I now know the final answer\n"
    "Final Answer: ..."
)
REACTTRACE = (
    "Thought: I need the authors of the ReAct paper first.\n"
    "Action: web_search\n"
    "Action Input: ReAct paper LLM reasoning acting authors\n"
    "Observation: Yao et al., ICLR 2023 (arXiv:2210.03629).\n"
    "\n"
    "Thought: Now the year of the transformer paper.\n"
    "Action: web_search\n"
    "Action Input: transformer paper Attention Is All You Need year\n"
    "Observation: Published in 2017 (Vaswani et al.).\n"
    "\n"
    "Thought: I have both. 2023 - 2017.\n"
    "Action: calculator\n"
    "Action Input: 2023 - 2017\n"
    "Observation: 6\n"
    "\n"
    "Thought: I now know the final answer.\n"
    "Final Answer: ReAct (Yao et al., 2023) came six years after\n"
    "the transformer paper (Vaswani et al., 2017)."
)
RAGDIAG = (
    "  PDFs -> text -> chunks -> embeddings -> vector store\n"
    "                                              |\n"
    "                                              v\n"
    "  question -> embed -> top-k chunks -> prompt + chunks -> LLM -> answer\n"
    "                                                                  |\n"
    "                                                                  v\n"
    "                                                        [chunk-id citations]"
)
W3NODE = (
    "                +------------------------------------------+\n"
    "   user query   |             LangGraph orchestrator       |\n"
    "       |        |                                          |\n"
    "       +------> |  +---------+   +----------+   +---------+ |\n"
    "                |  | Planner |-->|Retriever |-->| Writer  | |\n"
    "                |  +---------+   +----------+   +---------+ |\n"
    "                |                      ^                    |\n"
    "                |            hybrid_retriever()             |\n"
    "                |            ^ straight from Notebook 03     |\n"
    "                +------------------------------------------+"
)
STRUCT_CODE = (
    'from pydantic import BaseModel, Field\n'
    'from typing import Literal\n'
    '\n'
    'class ClaimCheck(BaseModel):\n'
    '    verdict: Literal["true", "false", "uncertain"]\n'
    '    confidence: Literal["low", "medium", "high"]\n'
    '    reason: str = Field(description="<= 15 words")\n'
    '\n'
    'checker = llm.with_structured_output(ClaimCheck)\n'
    'result = checker.invoke("The transformer paper is from 2017. True?")\n'
    '\n'
    'result.verdict      # -> "true"  (a typed object, not a string)'
)

DECK = [
    {"type": "title", "title": "Agentic AI Workshop",
     "lede": "Week 2 — From LLMs to Agents: why language models work, and the ladder that turns one into an agent.",
     "byline": "Dr. Kerem Delikoyun  ·  TUMCREATE  ·  3 June 2026", "building": True},

    {"type": "opener", "chapter": "Block 0 · 5 min", "title": "Where we left off",
     "blurb": "Last week: how we got here, plus Claude Code and Organon — working agents, built as plain files on disk. Today we open the hood."},

    {"type": "content", "title": "Week 1 recap, in one slide", "body": [
        ("cols",
         [("h", "What Week 1 covered"),
          ("b", ["**Foundations** — two long paths (a 200-yr abstraction staircase + an 80-yr AI road) converge at natural language",
                 "**Claude Code** — the CLI-native coding agent + its primitives: `CLAUDE.md`, tools, skills, sub-agents, MCP, hooks",
                 "**Organon** — an agentic OS for scientists; memory, identity & skills as files on disk",
                 "The equation: **agent = LLM + primitives** (memory, identity, tools, learnings)"])],
         [("h", "Where Week 2 goes"),
          ("n", ["*Open the box* — why LLMs work at all",
                 "*Climb the ladder* — bare model -> tools -> reasoning",
                 "*Pick up a framework* — LangChain",
                 "*Give it memory* — RAG"])]),
    ], "foot": ("muted", "Five Colab notebooks back every idea on the right. You'll run them as we go.")},

    {"type": "opener", "chapter": "Part 1 · Why LLMs work", "title": "Opening the box",
     "blurb": "Six slides on what a language model actually is — so \"agent\" isn't magic stacked on magic."},

    {"type": "content", "title": "The unlock: one architecture, three exponentials", "body": [
        ("cols",
         [("p", "For 60 years, language AI was hand-built rules and narrow models. Three things changed at once:"),
          ("b", ["**Architecture** — the Transformer (2017): attention replaces recurrence, so training parallelises across a whole sequence",
                 "**Scale** — parameters from millions -> hundreds of billions",
                 "**Data + compute** — train on much of the public web"])],
         [("card", "The bitter lesson", "blue", "General methods that scale with compute beat clever hand-crafted ones. The Transformer was the first language architecture that kept getting better the more data and compute you fed it — no plateau in sight."),
          ("muted", "\"Attention Is All You Need\", Vaswani et al. 2017 — one of the five papers in today's RAG corpus.")]),
    ]},

    {"type": "content", "title": "Under the hood: the Transformer (2017)", "wide": True,
     "image": (TXFIG, 9.15, 1.40, 3.45), "body": [
        ("cols",
         [("p", "Every modern LLM is built from this. Two stacks of identical blocks; each block is just three ideas:"),
          ("b", ["**Multi-Head Attention** — every token weighs what's relevant in all the others (the engine from the last slide)",
                 "**Feed-Forward** — a small per-token network on top",
                 "**Add & Norm** — residuals + normalisation that keep deep stacks trainable"]),
          ("p", "Stack it **N×**, add embeddings + positional encoding — that's the model. Chat LLMs keep the right (decoder) stack."),
          ("card", "Why it mattered", "blue", "No recurrence: the whole sequence is processed in **parallel** — what made web-scale training possible, and the **T** in GPT."),
          ("muted", "Figure 1, \"Attention Is All You Need\", Vaswani et al. 2017 — paper #2 in today's RAG corpus.")],
         [])  # right column intentionally empty; the figure sits here
    ]},

    {"type": "content", "title": "One dumb objective: predict the next token", "wide": True, "body": [
        ("cols",
         [("dgm", "nexttok"),
          ("p", "Train on trillions of tokens with a single self-supervised target: **guess the next token**. No labels, no human annotation — the text is its own answer key.")],
         [("h", "Why this is enough."),
          ("p", "To predict the next token well across *all* of human text, the model is forced to learn grammar, facts, translation, arithmetic, code, and a little reasoning."),
          ("card", "Compression ≈ understanding", "blue", "Predicting well means modelling the structure that *generated* the text. Squeezing the web into weights forces useful abstractions to fall out.")]),
    ]},

    {"type": "content", "title": "How text becomes math: tokens -> embeddings", "body": [
        ("cols",
         [("p", "**Tokenisation** — text is split into sub-word *tokens*."),
          ("code", '"agentic"  ->  ["agent", "ic"]\n"RAG"      ->  ["R", "AG"]'),
          ("p", "**Embedding** — each token maps to a vector in a high-dimensional space (hundreds–thousands of dims).")],
         [("p", "Meaning becomes *geometry*: tokens with similar meaning sit close together."),
          ("code", "king  - man  + woman  ~  queen\nParis - France + Japan ~ Tokyo"),
          ("card", "Why it matters today", "blue", "The **same** embedding idea powers RAG retrieval in Notebook 03 — we embed text chunks and find the nearest ones to a question.")]),
    ]},

    {"type": "content", "title": "The engine: attention as learned relevance", "wide": True, "body": [
        ("cols",
         [("dgm", "attn"),
          ("p", "Each token builds its meaning by **attending** to the other tokens that matter — weighting them by learned relevance. Stack this dozens of times and you get rich, context-aware representations.")],
         [("h", "Two properties that changed everything:"),
          ("b", ["**Parallel** — the whole sequence is processed at once (RNNs went word-by-word). This is what made web-scale training feasible.",
                 "**Long-range** — a token can attend to anything in the context window, not just its neighbours."]),
          ("muted", "A Transformer block = attention + a small feed-forward network, repeated N times.")]),
    ]},

    {"type": "content", "title": "Scaling laws & emergence", "body": [
        ("cols",
         [("p", "Performance improves *predictably* with model size, data, and compute — smooth power-law curves (Kaplan 2020; Hoffmann/**Chinchilla** 2022)."),
          ("card", "Chinchilla, compute-optimal", "blue", "For a fixed compute budget, scale parameters and data *together* — about **20 training tokens per parameter**. Most early models were badly under-trained."),
          ("muted", "Chinchilla is paper #3 in today's RAG corpus — you'll retrieve this exact claim in Notebook 03.")],
         [("p", "**Emergence** — some abilities are absent in small models and appear, fairly abruptly, past a scale threshold:"),
          ("b", ["in-context learning (learn from examples in the prompt, no weight update)",
                 "multi-step arithmetic",
                 "following instructions",
                 "**chain-of-thought reasoning**"]),
          ("p", "The capabilities we build on this week *emerged* from scale — nobody hand-coded them.")]),
    ]},

    {"type": "content", "title": "Why this architecture fits us so well", "body": [
        ("cols",
         [("p", "**Language is humanity's API.** We already encode knowledge, intent, and reasoning in text. A model fluent in text plugs straight into how we work — no new interface to learn."),
          ("p", "**In-context learning ≈ working memory.** Put examples or facts in the prompt and behaviour adapts on the spot, like a person holding instructions in mind for one task.")],
         [("card", "System 1 vs System 2 (Kahneman)", "green", "**System 1** — fast, automatic, intuitive. A raw next-token answer is System-1-like: one fluent pass, no deliberation.||**System 2** — slow, effortful, step-by-step. We *elicit* it by asking the model to reason out loud."),
          ("p", "That single idea — **make the fast model slow down and deliberate** — is the seed of every reasoning pattern in Part 2.")]),
    ], "foot": ("muted", "Analogy, not a claim that models literally think. It's a useful map for *why* the patterns work.")},

    {"type": "opener", "chapter": "Part 2 · Single agents", "title": "The agentic ladder",
     "blurb": "From a bare model to a reasoning, tool-using agent — one rung at a time. All single-agent this week; multi-agent is Week 3."},

    {"type": "content", "title": "The ladder — what we climb today", "body": [
        ("full", [("dgm", "ladder")]),
        ("cols",
         [("p", "Each rung **adds one capability** to the rung below. Nothing here is a new model — it's the same LLM with more scaffolding around it.")],
         [("card", "The thread", "blue", "L0 answers. L1 lets it *act*. L2 lets it *deliberate*. The \"agent\" is the loop you wrap around the model — not the model itself.")]),
    ]},

    {"type": "content", "title": "Anatomy of an agent", "body": [
        ("full", [("dgm", "anatomy"),
                  ("p", "**Tools** = capability. **Memory** = state. **Loop** = agency. Take away the loop and you have an LLM call. Take away tools and you have a chatbot. Take away memory and you have amnesia.")]),
    ]},

    {"type": "content", "title": "L0 — the bare LLM", "body": [
        ("cols",
         [("p", "One prompt in, one answer out. Brilliant at language; blind beyond its training."),
          ("h", "What it cannot do:"),
          ("b", ["Know anything after its training cut-off",
                 "Look up a fact it didn't memorise",
                 "Do anything in the world (send mail, query a DB, run code)",
                 "Reliably do long exact arithmetic",
                 "Tell you *why* it's confident"])],
         [("code", "user -> LLM -> answer"),
          ("card", "The ceiling", "coral", "A bare LLM is a brilliant *improviser* with no hands, no memory of today, and no scratchpad. Everything above this rung removes one of those limits.")]),
    ]},

    {"type": "content", "title": "L1 — tool use: give the model hands", "wide": True, "body": [
        ("cols",
         [("p", "A **tool** is just a function the model may ask you to call. The loop:"),
          ("dgm", "l1loop"),
          ("p", "The model never runs anything. It *requests*; your code executes and returns the result.")],
         [("p", "**Function calling.** Modern models are fine-tuned to pick a tool and fill its arguments from a schema you provide."),
          ("card", "Who carries the responsibility?", "blue", "The **vendor** trained the tool-selection skill. **You** own writing clear tool names + descriptions. Get those right and the model routes well."),
          ("muted", "Suddenly the model can fetch fresh facts, query data, and act — L0's ceiling is gone.")]),
    ]},

    {"type": "content", "title": "L2 — reasoning: make the fast model slow down", "body": [
        ("full", [("p", "The jump from L1 to L2 is **System 1 -> System 2**. Same model, but we change the prompt so it *deliberates* before it answers or acts.")]),
        ("cols",
         [("card", "Chain-of-Thought", "blue", "Think in a straight line, out loud."),
          ("card", "Tree-of-Thoughts", "blue", "Branch, evaluate, search — don't bet on the first idea.")],
         [("card", "ReAct", "green", "Reason *and* act — interleave thinking with tool calls."),
          ("muted", "CoT is the base. ToT searches over it. ReAct couples it to L1's tools. The next four slides take each in turn.")]),
    ]},

    {"type": "content", "title": "Chain-of-Thought (CoT)", "wide": True, "body": [
        ("cols",
         [("p", "Five words change the answer: **\"Let's think step by step.\"**"),
          ("code", COTBLK)],
         [("p", "The model uses its own intermediate tokens as a **scratchpad** — deliberate reasoning, summoned by a prompt."),
          ("card", "Why it works", "blue", "Hard problems need intermediate steps. Forcing the model to write them spends compute *before* the answer, instead of guessing in one pass."),
          ("muted", "CoT is System 2, on demand. It's the building block ReAct extends.")]),
    ]},

    {"type": "content", "title": "Tree-of-Thoughts (ToT)", "wide": True, "body": [
        ("cols",
         [("p", "CoT commits to one line of reasoning. If an early step is wrong, the whole chain is wrong. **ToT treats reasoning as search.**"),
          ("dgm", "tot")],
         [("p", "**Propose -> evaluate -> expand the best**, with backtracking out of dead ends."),
          ("card", "The trade", "blue", "ToT spends a lot more compute (many model calls) to buy accuracy on problems where one CoT chain routinely fails — planning, puzzles, proof search."),
          ("muted", "Notebook 02 shows one propose/evaluate level live on the Game of 24.")]),
    ]},

    {"type": "content", "title": "ReAct = Reason + Act", "wide": True, "body": [
        ("cols",
         [("p", "CoT thinks. L1 acts. **ReAct interleaves them** in one loop (Yao et al. 2022): a Thought, an Action (tool call), the Observation it returns — repeat."),
          ("code", REACTPROMPT)],
         [("p", "**There is no \"ReAct API\".** It's English in a fixed format that the model produces and *we* parse."),
          ("card", "The mechanism", "blue", "Stop the model right after \"Action Input\". Run the tool yourself. Append the real \"Observation\". Re-prompt. The model never hallucinates a result — it reads the one your code returned.")]),
    ]},

    {"type": "content", "title": "What a ReAct trace looks like", "body": [
        ("full", [("code", REACTTRACE),
                  ("muted", "Thought -> Action -> Observation, repeat. That loop *is* the agent. You'll build it by hand in Notebook 02.")]),
    ]},

    {"type": "content", "title": "Two ways to pick a tool", "body": [
        ("cols",
         [("h3", "ReAct agent"),
          ("p", "Tool choice lives in the **prompt** you write. The model emits `Action: / Action Input:` as text; you parse it with a regex."),
          ("b", ["**High control** — read and edit every step",
                 "Brittle parsing; works even on models with no native tool API"])],
         [("h3", "Function-calling agent"),
          ("p", "Tool choice is the **model vendor's** native capability. It returns a structured tool call; no parsing."),
          ("b", ["**Less control**, much less headache",
                 "Robust; needs a modern tool-calling model"])]),
        ("full", [("card", "\"Where do we shift the responsibility?\"", "blue",
                   "ReAct keeps tool selection in *your* prompt; the function-calling agent shifts it to the *vendor*. Most production agents use function calling; ReAct remains the clearest way to *understand* what an agent is.")]),
    ]},

    {"type": "opener", "chapter": "Part 3 · The framework", "title": "A short intro to LangChain",
     "blurb": "We've climbed the ladder by hand. LangChain is the standard toolkit that implements every rung — so you don't re-write the loop each time."},

    {"type": "content", "title": "What LangChain is (and isn't)", "body": [
        ("cols",
         [("p", "A thin, **vendor-neutral** layer over \"call a model, give it a prompt, maybe tools, parse the output.\" The same code runs on Anthropic, OpenAI, Google."),
          ("h", "Five pieces, one operator:"),
          ("table", ["Piece", "Class"],
           [["Model wrapper", "`ChatAnthropic`"], ["Messages", "`SystemMessage` / `HumanMessage`"],
            ["Prompt template", "`ChatPromptTemplate`"], ["Output parser", "`StrOutputParser`"],
            ["Tool", "`@tool`"]])],
         [("p", "**LCEL — the pipe.** Connect components with `|`, like a Unix pipe:"),
          ("code", 'chain = prompt | llm | parser\nchain.invoke({"concept": "embeddings"})'),
          ("card", "It isn't magic", "blue", "Everything we built by hand — the tool loop, the ReAct parse — LangChain wraps in one call. Because you built the primitive first, the abstraction is transparent.")]),
    ]},

    {"type": "content", "title": "LangChain implements the ladder", "body": [
        ("full", [
            ("table", ["Rung", "By hand (Parts 1-2)", "In LangChain"],
             [["**L0** bare model", "`client.messages.create(...)`", "`ChatAnthropic(...).invoke(...)`"],
              ["**L1** tool use", "your `for` loop + regex", "`create_tool_calling_agent` + `AgentExecutor`"],
              ["**L2** ReAct", "your prompt + `parse_action`", "`create_react_agent` + `AgentExecutor`"],
              ["**Memory** (RAG)", "manual chunk + embed + search", "loaders + splitters + vector stores + retrievers"]]),
            ("muted", "We pin the classic LangChain 0.3.x line — the last with `AgentExecutor`. LangGraph (the newer agent runtime) is the Week 3 reveal."),
        ]),
    ]},

    {"type": "content", "title": "Structured output: function calling, aimed at the answer", "wide": True, "body": [
        ("cols",
         [("code", STRUCT_CODE)],
         [("p", "Point a schema at the **answer** instead of a tool, and the model fills it in. You get back a **validated Python object** — no regex, no \"please respond in JSON\"."),
          ("card", "Same machinery", "blue", "It's the exact function-calling mechanism from the tool-calling agent. Reach for it whenever an agent's output feeds the next pipeline step and must have predictable fields."),
          ("muted", "Run live in Notebook 01.")]),
    ]},

    {"type": "content", "title": "From demo to production", "body": [
        ("cols",
         [("card", "Real tools — Tavily", "blue",
           "Our notebooks mock `web_search` so they run offline. In production, drop in a hosted search built for agents — same tool interface, real results:"),
          ("code", "search = TavilySearchResults(max_results=3)\nagent  = create_react_agent(llm, [search], prompt)"),
          ("muted", "`langchain_community.tools.tavily_search`")],
         [("card", "Observability — LangSmith", "green",
           "Set two env vars and every agent run is traced step-by-step — each Thought, Action, tool latency, token count — with no code change:"),
          ("code", 'LANGCHAIN_TRACING_V2 = "true"\nLANGCHAIN_API_KEY    = "ls-..."'),
          ("p", "You debug an agent by **watching the trace**, not guessing.")]),
    ]},

    {"type": "content", "title": "Hands-on: three notebooks, back to back", "body": [
        ("cols",
         [("card", "Notebook 00 · LangChain in 15 min", "blue", "Model wrapper, messages, prompt templates, the LCEL pipe, first `@tool`."),
          ("card", "Notebook 01 · Tool use", "blue", "The smallest agent, three ways: hand-rolled -> native tool-use -> LangChain tool-calling agent.")],
         [("card", "Notebook 02 · ReAct, CoT, ToT", "green", "CoT -> hand-built ReAct loop -> `create_react_agent` -> ToT -> the function-calling-vs-ReAct contrast."),
          ("muted", "Same model. Same loop you just saw on slides. Now you watch it run, step by step.")]),
        ("full", [("run", "▶ Switch to Colab — Notebooks 00, 01, 02")]),
    ]},

    {"type": "opener", "chapter": "Part 4 · Memory", "title": "RAG end-to-end",
     "blurb": "Retrieval-Augmented Generation gives the agent memory it can cite. Two parts — both visible — compared three ways."},

    {"type": "content", "title": "RAG in one diagram", "body": [
        ("full", [("dgm", "ragdiag"),
                  ("p", "Two halves. **Retrieval** finds the relevant chunks (same embedding idea from Part 1). **Augmentation** stuffs them into the prompt. The model's job becomes \"answer from this context, and cite it.\"")]),
    ]},

    {"type": "content", "title": "Build it with LangChain pieces", "body": [
        ("cols",
         [("p", "Each pipeline stage is one LangChain class:"),
          ("table", ["Stage", "Class"],
           [["Load PDFs", "`PyPDFLoader`"], ["Split", "`RecursiveCharacterTextSplitter`"],
            ["Embed", "`HuggingFaceEmbeddings`"], ["Store (dense)", "`FAISS` · `Chroma`"],
            ["Store (sparse)", "`BM25Retriever`"], ["Fuse", "`EnsembleRetriever`"]])],
         [("p", "**FAISS vs Chroma** — both local, both `.as_retriever()`:"),
          ("card", "FAISS", "blue", "Fastest pure-vector index, lives in RAM. Great default."),
          ("card", "Chroma", "green", "A small persistent local DB — write once, reopen later, no re-embedding."),
          ("muted", "Same retriever interface — the rest of the pipeline doesn't care which you pick.")]),
    ]},

    {"type": "content", "title": "Three retrieval strategies", "body": [
        ("cols",
         [("p", "**Dense (FAISS / Chroma)**"),
          ("b", ["Embed query + chunks, cosine similarity", "Great for paraphrase & synonyms", "\"What does X *mean*?\""]),
          ("p", "**Sparse (BM25)**"),
          ("b", ["Classical keyword scoring, no embeddings", "Great for exact terms", "\"What about *RLAIF*?\""])],
         [("p", "**Hybrid (EnsembleRetriever)**"),
          ("b", ["Fuses both with Reciprocal Rank Fusion", "Catches paraphrase **and** exact-term", "Usually wins on section-level hit"]),
          ("card", "Notebook 04 appendix", "green", "The same recipe on real Elasticsearch with native RRF — production substrate, same idea.")]),
    ]},

    {"type": "content", "title": "The dirty secret of RAG: chunking dominates", "body": [
        ("cols",
         [("p", "**Fixed-window** — split every N chars"),
          ("b", ["Fast; breaks sentences mid-thought"]),
          ("p", "**Recursive** — paragraph -> sentence -> word"),
          ("b", ["Respects boundaries; LangChain default", "Better for dense *and* sparse"]),
          ("p", "**Semantic** — split on embedding-distance jumps"),
          ("b", ["Slowest, sometimes best for narrative; we mention, don't build"])],
         [("card", "Notebook 03 · chunk + overlap", "blue", "137 pages -> **676 chunks** at 800 chars / 100 overlap.||Chunk size and overlap usually move retrieval quality *more* than the choice of embedding model. Sweep them on your own corpus."),
          ("p", "For the workshop we use **recursive** — the cheapest win.")]),
    ]},

    {"type": "content", "title": "Comparison: hit@5 on five questions", "body": [
        ("cols",
         [("p", "**Test set** — 5 questions, each labelled with its ground-truth paper + section keywords. Two metrics:"),
          ("b", ["**paper-hit@5** — right *paper* in the top 5?", "**section-hit@5** — right *part* of it?"]),
          ("p", "The first is easy; the second is honest.")],
         [("p", "**Notebook 03 run (5-paper corpus):**"),
          ("table", ["", "paper-hit", "section-hit"],
           [["Dense", "5/5", "5/5"], ["BM25", "5/5", "5/5"], ["Hybrid", "5/5", "**5/5**"]]),
          ("p", "On a small, clean corpus all three nail it. On a big messy one, hybrid pulls ahead — it catches both query types.")]),
    ]},

    {"type": "content", "title": "RAG vs raw LLM — same question", "body": [
        ("cols",
         [("card", "Raw LLM", "coral", "\"The ReAct loop interleaves reasoning and action...\"||*No citation. Possibly wrong on specifics. Sounds confident either way.*")],
         [("card", "RAG", "green", "\"...interleaving reasoning, action and observation improves groundedness **[react_yao_2022.pdf#540]**.\"||*Every claim points to a chunk you can re-read.*")]),
        ("full", [("p", "**Citations are the killer feature** — not prompt-engineering magic, just retrieval plus telling the model to cite its source."),
                  ("run", "▶ Switch to Colab — Notebook 03 (run the comparison, then ask a grounded question)")]),
    ]},

    {"type": "opener", "chapter": "Part 5 · 10 min", "title": "What's next — Week 3",
     "blurb": "A single ReAct + RAG agent covers a lot. Specialisation is the honest reason to reach for more than one."},

    {"type": "content", "title": "Limits of the single-agent pattern", "body": [
        ("full", [
            ("table", ["Limit", "Symptom", "Fix in W3"],
             [["Context overflow", "Long reasoning traces exhaust the window", "Hand subtasks to specialists"],
              ["No specialisation", "One prompt does everything, none of it well", "Planner / Retriever / Writer split"],
              ["Sequential bottleneck", "Parallel-decomposable work runs serially", "Graph orchestration"],
              ["Hard to evaluate", "Can't isolate which step failed", "Per-node metrics"]]),
            ("muted", "These are the **L3** rung — multi-agent systems — that we deliberately left off the ladder today."),
        ]),
    ]},

    {"type": "content", "title": "Today's retriever becomes a W3 tool node", "body": [
        ("full", [("dgm", "w3node"),
                  ("p", "The retrieval *function* doesn't change — the orchestrator around it does. Same with your ReAct loop: it becomes a node in a graph.")]),
    ]},

    {"type": "content", "title": "What to read before Week 3", "body": [
        ("full", [("b", ["**LangGraph quickstart** — `langchain-ai/langgraph` repo",
                         "**Anthropic's \"Building effective agents\"** — the workflows-vs-agents section",
                         "**GPT-Researcher** — a production multi-agent deep-research pattern",
                         "Your **Notebook 02 ReAct trace** — bring it back; we'll convert it to a graph"])]),
    ]},

    {"type": "opener", "chapter": "Wrap-up", "title": "Take it home",
     "blurb": "Five notebooks, one mental model: the loop is the agent."},

    {"type": "content", "title": "What you take home", "body": [
        ("full", [("b", ["**Notebook 00** — LangChain in 15 minutes",
                         "**Notebook 01** — tool use, three ways (hand-rolled -> native -> tool-calling agent)",
                         "**Notebook 02** — CoT, ReAct from scratch & via `create_react_agent`, ToT, FC-vs-ReAct",
                         "**Notebook 03** — full RAG with FAISS + Chroma + hybrid retrieval",
                         "**Notebook 04** — optional Elasticsearch appendix"]),
                  ("p", "All open in Colab from **github.com/krmdel/agentic-ai-for-scientists-workshop/week-02-patterns**")]),
    ]},

    {"type": "content", "title": "One sentence to remember", "body": [
        ("full", [("card", "The Week 2 takeaway", "green",
                   "A language model predicts the next token. An **agent** is the loop you wrap around it — tools give it hands, reasoning gives it a scratchpad, RAG gives it memory it can cite. Frameworks add ergonomics; they don't add the agency.")]),
    ]},

    {"type": "title", "title": "Questions?",
     "lede": "Kerem Delikoyun  ·  TUMCREATE",
     "byline": "kerem.delikoyun@tum-create.edu.sg  ·  github.com/krmdel/agentic-ai-for-scientists-workshop", "building": False},
]


def main():
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH
    for i, sl in enumerate(DECK):
        page = i + 1
        t = sl["type"]
        if t == "title":
            build_title(prs, sl["title"], sl["lede"], sl["byline"], sl.get("building", False))
        elif t == "opener":
            build_opener(prs, page, sl["chapter"], sl["title"], sl["blurb"])
        else:
            build_content(prs, page, sl)
    prs.save(str(OUT))
    print(f"wrote {OUT}  ({len(DECK)} slides)")


if __name__ == "__main__":
    main()
